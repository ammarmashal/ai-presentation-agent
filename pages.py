'''def create_title_slide(prs, presentation_title, text_rgb):
    """Create a title slide with theme colors"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = presentation_title
    subtitle.text = "Created with AI Presentation Generator"
    
    # Apply text color with bold for better visibility
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.color.rgb = RGBColor(*text_rgb)
        paragraph.font.bold = True
        paragraph.font.size = Pt(44)
    
    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.font.color.rgb = RGBColor(*text_rgb)
        paragraph.font.size = Pt(18)
    
    return slide'''

from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE, MSO_FILL_TYPE, MSO_FILL
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image as PILImage
import io

def create_image_title_slide(prs, presentation_title, image_path, primary_rgb, text_rgb, accent_rgb):
    """
    Creates a title slide with an image on the left and title text on the right,
    mimicking a modern, stylish layout.
    """
    # Use a blank slide layout (index 6) for a custom layout
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Get slide dimensions for proportional scaling
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # --- Set Background for the Right Half (Gradient) ---
    right_half_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        slide_width / 2, 
        Inches(0), 
        slide_width / 2, 
        slide_height
    )
    fill = right_half_shape.fill

    # Initialize gradient fill
    fill.gradient()
    
    # Set gradient properties
    fill.gradient_angle = 180

    # Handle gradient stops
    start_color_gradient = RGBColor(255, 255, 255)  # White at the top
    end_color_gradient = RGBColor(*primary_rgb)
    
    # Ensure we have at least 2 gradient stops
    while len(fill.gradient_stops) < 2:
        fill.gradient_stops.add_stop(0.0, RGBColor(0, 0, 0))
    
    # Modify the first two stops
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[0].color.rgb = start_color_gradient
    
    fill.gradient_stops[1].position = 1.0
    fill.gradient_stops[1].color.rgb = end_color_gradient
    
    # Remove any additional stops if they exist
    for i in range(len(fill.gradient_stops) - 1, 1, -1):
        fill.gradient_stops[i].delete()
    
    right_half_shape.line.fill.background()

    # --- Add Image on the Left ---
    img_left = Inches(0)
    img_top = Inches(0)
    img_width = slide_width / 2
    img_height = slide_height

    try:
        with PILImage.open(image_path) as img:
            original_width_px, original_height_px = img.size
        
        container_aspect_ratio = float(img_width) / float(img_height)
        image_aspect_ratio = float(original_width_px) / float(original_height_px)

        if image_aspect_ratio > container_aspect_ratio:
            new_height_emu = img_height
            new_width_emu = int(img_height * image_aspect_ratio)
            picture_left_emu = img_left.emu + (img_width - new_width_emu) / 2
            picture_top_emu = img_top.emu
        else:
            new_width_emu = img_width
            new_height_emu = int(img_width / image_aspect_ratio)
            picture_left_emu = img_left.emu
            picture_top_emu = img_top.emu + (img_height - new_height_emu) / 2

        slide.shapes.add_picture(
            image_path, picture_left_emu, picture_top_emu,
            width=new_width_emu, height=new_height_emu
        )
    
    except FileNotFoundError:
        print(f"Error: Image file not found at {image_path}. Using a grey placeholder.")
        placeholder_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, img_left, img_top, img_width, img_height
        )
        placeholder_fill = placeholder_shape.fill
        placeholder_fill.solid()
        placeholder_fill.fore_color.rgb = RGBColor(180, 180, 180)
        placeholder_shape.line.fill.background()
    
    # --- Add Title and Subtitle Text on the Right ---
    title_box_left = slide_width / 2 + Inches(0.5)
    title_box_top = Inches(1.5)
    title_box_width = slide_width / 2 - Inches(1)
    title_box_height = Inches(3)

    title_box = slide.shapes.add_textbox(title_box_left, title_box_top, title_box_width, title_box_height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    
    # CORRECTED: Simple and reliable way to clear and add text
    # Clear all text first
    title_frame.clear()
    
    # Add title paragraph
    p_title = title_frame.paragraphs[0]
    p_title.text = presentation_title
    p_title.alignment = PP_ALIGN.CENTER
    p_title.font.size = Pt(44)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(*accent_rgb)

    # Add subtitle paragraph
    p_subtitle = title_frame.add_paragraph()
    p_subtitle.text = "Welcome to an exploration of its pivotal milestones, and its pervasive impact on our daily lives and industries."
    p_subtitle.alignment = PP_ALIGN.CENTER
    p_subtitle.font.size = Pt(18)
    p_subtitle.font.color.rgb = RGBColor(*text_rgb)
    p_subtitle.space_before = Inches(0.3)

    title_box.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return slide
# And a simple fallback create_title_slide (if not already defined and used in your main loop)
def create_title_slide(prs, presentation_title, text_rgb):
    """
    Creates a basic text-only title slide (fallback for image issues).
    """
    slide = prs.slides.add_slide(prs.slide_layouts[0]) # Use default title layout
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = presentation_title
    subtitle.text = "Created with AI Presentation Generator"
    
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.color.rgb = RGBColor(*text_rgb)
        paragraph.font.bold = True
        paragraph.font.size = Pt(44)
    
    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.font.color.rgb = RGBColor(*text_rgb)
        paragraph.font.size = Pt(18)
    
    return slide

'''def create_image_title_slide(prs, presentation_title, image_path, text_rgb):
    """
    Creates a title slide with an image on the left and title text on the right,
    using a blank slide layout for custom positioning.
    """
    # Use a blank slide layout (index 6) for a custom layout
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Get slide dimensions for proportional scaling
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # --- Add Image on the Left ---
    # Position and size the image to take up the left half of the slide
    img_left = Inches(0)
    img_top = Inches(0)
    img_width = slide_width / 2
    img_height = slide_height

    # Load and scale image to fit container while maintaining aspect ratio
    try:
        with Image.open(image_path) as img:
            original_width_px, original_height_px = img.size
    except FileNotFoundError:
        print(f"Error: Image file not found at {image_path}. Skipping image.")
        # Fallback to a plain text title slide if the image isn't found
        return create_title_slide(prs, presentation_title, text_rgb)

    container_width_emu = img_width
    container_height_emu = img_height

    ratio_w = container_width_emu / original_width_px
    ratio_h = container_height_emu / original_height_px
    scale_factor = min(ratio_w, ratio_h)
    
    new_width_emu = int(original_width_px * scale_factor)
    new_height_emu = int(original_height_px * scale_factor)
    
    # Center the scaled picture within the left half of the slide
    picture_left_emu = img_left.emu + (container_width_emu - new_width_emu) / 2
    picture_top_emu = img_top.emu + (container_height_emu - new_height_emu) / 2
    
    slide.shapes.add_picture(
        image_path, picture_left_emu, picture_top_emu,
        width=new_width_emu, height=new_height_emu
    )
    
    # --- Add Title on the Right ---
    # Position and size the title to take up the right half of the slide
    title_left = slide_width / 2
    title_top = Inches(0)
    title_width = slide_width / 2
    title_height = slide_height
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    
    # Use a single paragraph for the title text
    p = title_frame.paragraphs[0]
    p.text = presentation_title
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*text_rgb)
    
    # Vertically center the text within its box
    title_box.vertical_anchor = MSO_ANCHOR.MIDDLE

    # --- Add Decorative Elements (optional) ---
    decor_right = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(0.1), Inches(0.4), Inches(0.4))
    decor_right.fill.solid()
    decor_right.fill.fore_color.rgb = RGBColor(*text_rgb)
    decor_right.line.fill.background()
    
    decor_left = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.1), Inches(6.5), Inches(0.3), Inches(0.3))
    decor_left.fill.solid()
    decor_left.fill.fore_color.rgb = RGBColor(*text_rgb)
    decor_left.line.fill.background()

    return slide'''

def create_title_content_slide(prs, section, points, text_rgb):
    """Create a slide with title and content with theme colors"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = section
    
    # Apply text color to title
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.font.color.rgb = RGBColor(*text_rgb)
    
    text_frame = content_shape.text_frame
    text_frame.clear()
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.word_wrap = True
    
    first = True
    for point in points:
        level = point.get("level", 0)
        
        if first:
            p = text_frame.paragraphs[0]
            first = False
        else:
            p = text_frame.add_paragraph()
        
        p.text = point["text"]
        p.level = level
        p.font.color.rgb = RGBColor(*text_rgb)
    
    return slide

def create_image_left_text_right_slide(prs, section, points, image_path, primary_rgb, text_rgb, accent_rgb):
    """Create a slide with image on left and text on right with gradient background like first slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Get slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # --- Set Gradient Background for Entire Slide ---
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height)
    fill = bg_shape.fill
    fill.gradient()
    fill.gradient_angle = 180
    
    # Set gradient stops (white to primary color)
    while len(fill.gradient_stops) < 2:
        fill.gradient_stops.add_stop(0.0, RGBColor(0, 0, 0))
    
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[0].color.rgb = RGBColor(255, 255, 255)  # White at top
    fill.gradient_stops[1].position = 1.0
    fill.gradient_stops[1].color.rgb = RGBColor(*primary_rgb)   # Primary color at bottom
    
    bg_shape.line.fill.background()
    
    # Define layout dimensions
    text_left = Inches(5.2)
    text_top = Inches(1.5)
    text_width = Inches(4.3)
    text_height = Inches(4.5)
    img_left = Inches(0.8)
    img_top = Inches(1.5)
    img_width = Inches(4)
    img_height = Inches(4.5)

    # Add text container with modern styling
    text_container = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, text_left, text_top, text_width, text_height
    )
    # Semi-transparent white background for text container
    text_container.fill.solid()
    text_container.fill.fore_color.rgb = RGBColor(255, 255, 255)
    text_container.fill.fore_color.alpha = 0.8  # 80% opacity
    
    text_container.line.width = Pt(2)
    text_container.line.color.rgb = RGBColor(*accent_rgb)
    text_container.line.dash_style = MSO_LINE.SOLID  # Changed to solid line

    # Add title with gradient background
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = section
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(*accent_rgb)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Modern decorative line - CORRECTED
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.1), Inches(8), Inches(0.02))
    line_fill = line.fill
    line_fill.solid()  # Use solid color for simplicity, or gradient if needed
    line_fill.fore_color.rgb = RGBColor(*accent_rgb)
    line.line.fill.background()

    # Add image with modern container
    img_container = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, img_left, img_top, img_width, img_height
    )
    # White background with subtle shadow effect
    img_container.fill.solid()
    img_container.fill.fore_color.rgb = RGBColor(255, 255, 255)
    img_container.line.width = Pt(3)
    img_container.line.color.rgb = RGBColor(*accent_rgb)
    
    try:
        with Image.open(image_path) as img:
            original_width_px, original_height_px = img.size
    except FileNotFoundError:
        print(f"Error: Image file not found at {image_path}")
        # Add placeholder with gradient
        placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, img_left, img_top, img_width, img_height)
        placeholder_fill = placeholder.fill
        placeholder_fill.solid()  # Use solid color for simplicity
        placeholder_fill.fore_color.rgb = RGBColor(200, 200, 200)
        placeholder.line.width = Pt(2)
        placeholder.line.color.rgb = RGBColor(*text_rgb)
    else:
        # Calculate dimensions to maintain aspect ratio
        container_width_emu = img_width.emu
        container_height_emu = img_height.emu
        
        ratio_w = container_width_emu / original_width_px
        ratio_h = container_height_emu / original_height_px
        
        scale_factor = min(ratio_w, ratio_h)
        
        new_width_emu = int(original_width_px * scale_factor)
        new_height_emu = int(original_height_px * scale_factor)

        # Center the picture within its container
        picture_left_emu = img_left.emu + (container_width_emu - new_width_emu) / 2
        picture_top_emu = img_top.emu + (container_height_emu - new_height_emu) / 2
        
        picture = slide.shapes.add_picture(
            image_path, picture_left_emu, picture_top_emu, 
            width=new_width_emu, height=new_height_emu
        )
    
    # Add text box with modern styling
    text_inner_left = text_left + Inches(0.2)
    text_inner_top = text_top + Inches(0.3)
    text_inner_width = text_width - Inches(0.4)
    text_inner_height = text_height - Inches(0.4)
    
    text_box = slide.shapes.add_textbox(
        text_inner_left, text_inner_top, text_inner_width, text_inner_height
    )
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    # Add modern bullet points
    first = True
    for i, point in enumerate(points):
        level = point.get("level", 0)
        if first:
            p = text_frame.paragraphs[0]
            first = False
        else:
            p = text_frame.add_paragraph()
        
        # Modern bullet point styling
        p.text = "● " + point["text"]  # Using circle bullet instead of dot
        p.level = level
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(*text_rgb)
        p.space_after = Inches(0.1)
        
        # Add accent color to first few words or important text
        if i == 0:  # Make first point stand out
            p.font.bold = True
            p.font.color.rgb = RGBColor(*accent_rgb)

    # Modern decorative elements - CORRECTED
    decor_right = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(0.1), Inches(0.4), Inches(0.4))
    decor_right_fill = decor_right.fill
    decor_right_fill.solid()  # Use solid color for simplicity
    decor_right_fill.fore_color.rgb = RGBColor(*accent_rgb)
    decor_right.line.fill.background()
    
    decor_left = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.1), Inches(6.5), Inches(0.3), Inches(0.3))
    decor_left_fill = decor_left.fill
    decor_left_fill.solid()  # Use solid color for simplicity
    decor_left_fill.fore_color.rgb = RGBColor(*primary_rgb)
    decor_left.line.fill.background()
    
    return slide

def create_image_right_text_left_slide(prs, section, points, image_path, primary_rgb, text_rgb, accent_rgb):
    """Create a slide with image on right and text on left with gradient background like first slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Get slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # --- Set Gradient Background for Entire Slide ---
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height)
    fill = bg_shape.fill
    fill.gradient()
    fill.gradient_angle = 180
    
    # Set gradient stops (white to primary color)
    while len(fill.gradient_stops) < 2:
        fill.gradient_stops.add_stop(0.0, RGBColor(0, 0, 0))
    
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[0].color.rgb = RGBColor(255, 255, 255)  # White at top
    fill.gradient_stops[1].position = 1.0
    fill.gradient_stops[1].color.rgb = RGBColor(*primary_rgb)   # Primary color at bottom
    
    bg_shape.line.fill.background()
    
    # Define layout dimensions (mirrored from left version)
    text_left = Inches(0.5)
    text_top = Inches(1.5)
    text_width = Inches(4.5)
    text_height = Inches(4.5)
    img_left = Inches(5.2)
    img_top = Inches(1.5)
    img_width = Inches(4.3)
    img_height = Inches(4.5)

    # Add text container with modern styling
    text_container = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, text_left, text_top, text_width, text_height
    )
    # Semi-transparent white background for text container
    text_container.fill.solid()
    text_container.fill.fore_color.rgb = RGBColor(255, 255, 255)
    text_container.fill.fore_color.alpha = 0.8  # 80% opacity
    
    text_container.line.width = Pt(2)
    text_container.line.color.rgb = RGBColor(*accent_rgb)
    text_container.line.dash_style = MSO_LINE.SOLID

    # Add title with left alignment
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = section
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(*accent_rgb)
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    # Add image container
    img_container = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, img_left, img_top, img_width, img_height
    )
    img_container.fill.solid()
    img_container.fill.fore_color.rgb = RGBColor(255, 255, 255)
    img_container.line.width = Pt(3)
    img_container.line.color.rgb = RGBColor(*accent_rgb)
    
    try:
        with Image.open(image_path) as img:
            original_width_px, original_height_px = img.size
    except FileNotFoundError:
        print(f"Error: Image file not found at {image_path}")
        # Add gradient placeholder
        placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, img_left, img_top, img_width, img_height)
        placeholder_fill = placeholder.fill
        placeholder_fill.solid()  # Use solid color for simplicity
        placeholder_fill.fore_color.rgb = RGBColor(200, 200, 200)
        placeholder.line.width = Pt(2)
        placeholder.line.color.rgb = RGBColor(*text_rgb)
    else:
        # Calculate dimensions to maintain aspect ratio
        container_width_emu = img_width.emu
        container_height_emu = img_height.emu
        
        ratio_w = container_width_emu / original_width_px
        ratio_h = container_height_emu / original_height_px
        
        scale_factor = min(ratio_w, ratio_h)
        
        new_width_emu = int(original_width_px * scale_factor)
        new_height_emu = int(original_height_px * scale_factor)

        # Center the picture within its container
        picture_left_emu = img_left.emu + (container_width_emu - new_width_emu) / 2
        picture_top_emu = img_top.emu + (container_height_emu - new_height_emu) / 2
        
        picture = slide.shapes.add_picture(
            image_path, picture_left_emu, picture_top_emu, 
            width=new_width_emu, height=new_height_emu
        )
    
    # Add text box
    text_inner_left = text_left + Inches(0.2)
    text_inner_top = text_top + Inches(0.3)
    text_inner_width = text_width - Inches(0.4)
    text_inner_height = text_height - Inches(0.4)
    
    text_box = slide.shapes.add_textbox(
        text_inner_left, text_inner_top, text_inner_width, text_inner_height
    )
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    # Add modern bullet points
    first = True
    for i, point in enumerate(points):
        level = point.get("level", 0)
        if first:
            p = text_frame.paragraphs[0]
            first = False
        else:
            p = text_frame.add_paragraph()
        
        p.text = "● " + point["text"]
        p.level = level
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(*text_rgb)
        p.space_after = Inches(0.1)
        
        if i == 0:
            p.font.bold = True
            p.font.color.rgb = RGBColor(*accent_rgb)

    # Modern decorative elements - CORRECTED
    decor_right = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(0.1), Inches(0.4), Inches(0.4))
    decor_right_fill = decor_right.fill
    decor_right_fill.solid()  # Use solid color for simplicity
    decor_right_fill.fore_color.rgb = RGBColor(*accent_rgb)
    decor_right.line.fill.background()
    
    decor_left = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.1), Inches(6.5), Inches(0.3), Inches(0.3))
    decor_left_fill = decor_left.fill
    decor_left_fill.solid()  # Use solid color for simplicity
    decor_left_fill.fore_color.rgb = RGBColor(*primary_rgb)
    decor_left.line.fill.background()
    
    return slide



from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_two_column_slide(prs, section, points, text_rgb):
    """
    Creates a two-column text slide with modern styling, including borders,
    bullet points, and decorative elements.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # --- Define Layout Dimensions ---
    left_text_left = Inches(0.5)
    left_text_top = Inches(1.5)
    left_text_width = Inches(4.5)
    left_text_height = Inches(5.5)
    
    right_text_left = Inches(5.2)
    right_text_top = Inches(1.5)
    right_text_width = Inches(4.3)
    right_text_height = Inches(5.5)

    # --- Add Title and Decorative Line ---
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = section
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(*text_rgb)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.1), Inches(8), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(*text_rgb)
    line.line.fill.background()
    
    # --- Split points into two columns ---
    mid_point = len(points) // 2
    left_points = points[:mid_point]
    right_points = points[mid_point:]
    
    # --- Add Left Text Column Container and Content ---
    left_text_container = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left_text_left, left_text_top, left_text_width, left_text_height
    )
    left_text_container.fill.background()
    left_text_container.line.width = Pt(1)
    left_text_container.line.color.rgb = RGBColor(*text_rgb)
    left_text_container.line.dash_style = MSO_LINE.DASH
    
    left_text_box = slide.shapes.add_textbox(left_text_left + Inches(0.2), left_text_top + Inches(0.5), 
                                             left_text_width - Inches(0.4), left_text_height - Inches(0.4))
    left_text_frame = left_text_box.text_frame
    left_text_frame.word_wrap = True
    
    first = True
    for i, point in enumerate(left_points):
        p = left_text_frame.paragraphs[0] if first else left_text_frame.add_paragraph()
        first = False
        p.text = "•  " + point["text"]
        p.level = point.get("level", 0)
        p.font.size = Pt(20)  # Increased font size
        p.font.color.rgb = RGBColor(*text_rgb)
        p.space_after = Inches(0.15) if i < len(left_points) - 1 else Inches(0.1)
    
    # --- Add Right Text Column Container and Content ---
    right_text_container = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, right_text_left, right_text_top, right_text_width, right_text_height
    )
    right_text_container.fill.background()
    right_text_container.line.width = Pt(1)
    right_text_container.line.color.rgb = RGBColor(*text_rgb)
    right_text_container.line.dash_style = MSO_LINE.DASH

    right_text_box = slide.shapes.add_textbox(right_text_left + Inches(0.2), right_text_top + Inches(0.5), 
                                              right_text_width - Inches(0.4), right_text_height - Inches(0.4))
    right_text_frame = right_text_box.text_frame
    right_text_frame.word_wrap = True
    
    first = True
    for i, point in enumerate(right_points):
        p = right_text_frame.paragraphs[0] if first else right_text_frame.add_paragraph()
        first = False
        p.text = "•  " + point["text"]
        p.level = point.get("level", 0)
        p.font.size = Pt(20)  # Increased font size
        p.font.color.rgb = RGBColor(*text_rgb)
        p.space_after = Inches(0.15) if i < len(right_points) - 1 else Inches(0.1)

    # --- Add Decorative Elements ---
    decor_right = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(0.1), Inches(0.4), Inches(0.4))
    decor_right.fill.solid()
    decor_right.fill.fore_color.rgb = RGBColor(*text_rgb)
    decor_right.line.fill.background()
    
    decor_left = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.1), Inches(6.5), Inches(0.3), Inches(0.3))
    decor_left.fill.solid()
    decor_left.fill.fore_color.rgb = RGBColor(*text_rgb)
    decor_left.line.fill.background()
    
    return slide


'''def create_image_left_text_right_slide(prs, section, points, image_path, text_rgb):
    """Create a slide with image on left and text on right with theme colors"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = section
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(*text_rgb)
    
    # Add image on left
    img_left = Inches(0.5)
    img_top = Inches(1.5)
    img_width = Inches(4.5)
    img_height = Inches(5.5)
    slide.shapes.add_picture(image_path, img_left, img_top, width=img_width, height=img_height)
    
    # Add text on right
    text_left = Inches(5.5)
    text_top = Inches(1.5)
    text_width = Inches(4.5)
    text_height = Inches(5.5)
    text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    first = True
    for point in points:
        level = point.get("level", 0)
        
        if first:
            p = text_frame.paragraphs[0]
            first = False
        else:
            p = text_frame.add_paragraph()
        
        p.text = point["text"]
        p.level = level
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(*text_rgb)
    
    return slide'''

'''def create_image_right_text_left_slide(prs, section, points, image_path, text_rgb):
    """Create a slide with image on right and text on left with theme colors"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = section
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(*text_rgb)
    
    # Add text on left
    text_left = Inches(0.5)
    text_top = Inches(1.5)
    text_width = Inches(4.5)
    text_height = Inches(5.5)
    text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    first = True
    for point in points:
        level = point.get("level", 0)
        
        if first:
            p = text_frame.paragraphs[0]
            first = False
        else:
            p = text_frame.add_paragraph()
        
        p.text = point["text"]
        p.level = level
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(*text_rgb)
    
    # Add image on right
    img_left = Inches(5.5)
    img_top = Inches(1.5)
    img_width = Inches(4.5)
    img_height = Inches(5.5)
    slide.shapes.add_picture(image_path, img_left, img_top, width=img_width, height=img_height)
    
    return slide'''


'''def create_two_column_slide(prs, section, points, text_rgb):
    """Create a slide with two columns of text with theme colors"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = section
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(*text_rgb)
    
    # Split points into two columns
    mid_point = len(points) // 2
    left_points = points[:mid_point]
    right_points = points[mid_point:]
    
    # Add left column
    left_text_left = Inches(0.5)
    left_text_top = Inches(1.5)
    left_text_width = Inches(4.5)
    left_text_height = Inches(5.5)
    left_text_box = slide.shapes.add_textbox(left_text_left, left_text_top, left_text_width, left_text_height)
    left_text_frame = left_text_box.text_frame
    left_text_frame.word_wrap = True
    
    first = True
    for point in left_points:
        level = point.get("level", 0)
        
        if first:
            p = left_text_frame.paragraphs[0]
            first = False
        else:
            p = left_text_frame.add_paragraph()
        
        p.text = point["text"]
        p.level = level
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(*text_rgb)
    
    # Add right column
    right_text_left = Inches(5.5)
    right_text_top = Inches(1.5)
    right_text_width = Inches(4.5)
    right_text_height = Inches(5.5)
    right_text_box = slide.shapes.add_textbox(right_text_left, right_text_top, right_text_width, right_text_height)
    right_text_frame = right_text_box.text_frame
    right_text_frame.word_wrap = True
    
    first = True
    for point in right_points:
        level = point.get("level", 0)
        
        if first:
            p = right_text_frame.paragraphs[0]
            first = False
        else:
            p = right_text_frame.add_paragraph()
        
        p.text = point["text"]
        p.level = level
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(*text_rgb)
    
    return slide
'''


def create_conclusion_slide(prs, section, points, text_rgb):
    """Create a conclusion slide with theme colors"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = section
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.color.rgb = RGBColor(*text_rgb)
    
    # Add content in the center
    content_left = Inches(1)
    content_top = Inches(2)
    content_width = Inches(8)
    content_height = Inches(4)
    content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Add key points
    for i, point in enumerate(points):
        if i < 3:  # Limit to 3 key points
            p = content_frame.add_paragraph() if i > 0 else content_frame.paragraphs[0]
            p.text = point["text"]
            p.font.size = Pt(20)
            p.level = 0
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = RGBColor(*text_rgb)
    
    return slide

def create_thank_you_slide(prs, primary_rgb, text_rgb, accent_rgb):
    """Create a Thank You slide with theme colors"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add background color using theme's primary color
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(*primary_rgb)
    background.line.fill.background()
    
    # Add Thank You text
    thank_you_left = Inches(1)
    thank_you_top = Inches(2.5)
    thank_you_width = Inches(8)
    thank_you_height = Inches(2)
    thank_you_box = slide.shapes.add_textbox(thank_you_left, thank_you_top, thank_you_width, thank_you_height)
    thank_you_frame = thank_you_box.text_frame
    thank_you_frame.text = "Thank You!"
    thank_you_frame.paragraphs[0].font.size = Pt(48)
    thank_you_frame.paragraphs[0].font.bold = True
    thank_you_frame.paragraphs[0].font.color.rgb = RGBColor(*text_rgb)
    thank_you_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add smaller subtitle
    subtitle_left = Inches(1)
    subtitle_top = Inches(5)
    subtitle_width = Inches(8)
    subtitle_height = Inches(1)
    subtitle_box = slide.shapes.add_textbox(subtitle_left, subtitle_top, subtitle_width, subtitle_height)
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Questions & Discussion"
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    # Use accent color for subtitle or lighter version of text color
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(*accent_rgb)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def list_available_themes(theme_folder="themes"):
    """List all available themes in the specified folder"""
    if not os.path.exists(theme_folder):
        os.makedirs(theme_folder)
        return []
    
    themes = [f for f in os.listdir(theme_folder) if f.endswith(".pptx")]
    return themes