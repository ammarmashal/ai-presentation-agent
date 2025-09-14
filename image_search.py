import os
import requests
from dotenv import load_dotenv
import json
import tempfile
from pptx.util import Inches
from pptx.enum.text import MSO_ANCHOR
import random
import re
from collections import Counter

# Load environment variables
load_dotenv()
UNSPLASH_ACCESS_KEY = os.getenv("UNSPLASH_ACCESS_KEY")

def search_images(query, count=5):
    """
    Search for images using Unsplash API
    """
    if not UNSPLASH_ACCESS_KEY:
        print("❌ Unsplash access key not found. Please add UNSPLASH_ACCESS_KEY to your .env file.")
        return []
    
    try:
        url = f"https://api.unsplash.com/search/photos"
        params = {
            "query": query,
            "per_page": count,
            "client_id": UNSPLASH_ACCESS_KEY
        }
        
        response = requests.get(url, params=params, timeout=30)
        response.raise_for_status()
        
        data = response.json()
        images = []
        
        for result in data.get("results", []):
            images.append({
                "id": result["id"],
                "url": result["urls"]["regular"],
                "description": result.get("description", result.get("alt_description", "")),
                "download_url": result["urls"]["regular"]
            })
        
        return images
    
    except Exception as e:
        print(f"❌ Error searching images: {e}")
        return []

def download_image(image_url, filename):
    """
    Download an image from a URL
    """
    try:
        response = requests.get(image_url, timeout=30)
        response.raise_for_status()
        
        with open(filename, 'wb') as f:
            f.write(response.content)
        
        return True
    except Exception as e:
        print(f"❌ Error downloading image: {e}")
        return False

def add_image_to_slide(slide, image_path, position="left_bottom", size=Inches(3)):
    """
    Add an image to a slide at the specified position
    """
    try:
        if position == "left_bottom":
            # Position the image at the left bottom corner
            left = Inches(6.5)
            top = Inches(3)
        else:
            # Default position (center)
            left = Inches(3)
            top = Inches(2)
        
        slide.shapes.add_picture(image_path, left, top, width=size, height=size)
        return True
    except Exception as e:
        print(f"❌ Error adding image to slide: {e}")
        return False

def create_image_slide(prs, image_path, title_text=""):
    """
    Create a dedicated slide for an image (for detailed presentations)
    """
    try:
        # Use a blank layout for image slides (no header)
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Get slide dimensions
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # Define image size (adjust as needed)
        image_width = Inches(6)  # Set the desired width of the image
        image_height = Inches(4)  # Set the desired height of the image

        # Calculate centered position
        left = (slide_width - image_width) / 2
        top = (slide_height - image_height) / 2

        # Add the image to the slide
        slide.shapes.add_picture(image_path, left, top, width=image_width, height=image_height)

        # Add a title if provided
        if title_text:
            title_box = slide.shapes.title
            if title_box:
                title_box.text = title_text

        return slide
    except Exception as e:
        print(f"❌ Error creating image slide: {e}")
        return None

def get_relevant_image_queries(presentation_title, slide_title, content, is_detailed=False):
    """
    Generate relevant search queries for images based on presentation content
    For detailed presentations, use the entire LLM content
    """
    queries = []
    
    if is_detailed:
        # For detailed presentations, use the entire content from LLM
        content_text = " ".join([point["text"] for point in content])
        all_text = f"{presentation_title} {slide_title} {content_text}"
        
        # Extract meaningful words (nouns, adjectives)
        words = re.findall(r'\b[a-zA-Z]{4,}\b', all_text.lower())
        
        # Remove common stop words
        stop_words = {"the", "and", "or", "but", "in", "on", "at", "to", "for", 
                     "of", "with", "by", "that", "this", "these", "those", "is",
                     "are", "was", "were", "be", "been", "being", "have", "has",
                     "had", "do", "does", "did", "will", "would", "could", "should"}
        
        meaningful_words = [word for word in words if word not in stop_words]
        
        # Get the most frequent words
        word_counts = Counter(meaningful_words)
        most_common = [word for word, count in word_counts.most_common(8)]
        
        queries.extend(most_common)
        
        # Add combinations
        for word in most_common[:3]:
            queries.append(f"{presentation_title} {word}")
            queries.append(f"{slide_title} {word}")
    else:
        # For simple presentations, use a more targeted approach
        content_text = " ".join([point["text"] for point in content])
        all_text = f"{slide_title} {content_text}"
        
        # Find specific nouns (more likely to have good images)
        words = re.findall(r'\b[a-zA-Z]{4,}\b', all_text.lower())
        
        # Filter out abstract terms
        abstract_terms = {"management", "strategy", "system", "process", "method", 
                         "approach", "concept", "theory", "principle", "framework",
                         "model", "analysis", "development", "implementation"}
        
        concrete_words = [word for word in words if word not in abstract_terms]
        
        # Add the most frequent concrete words
        word_counts = Counter(concrete_words)
        most_common = [word for word, count in word_counts.most_common(5)]
        
        queries.extend(most_common)
        
        # Add combinations with presentation title for context
        for word in most_common[:2]:
            queries.append(f"{presentation_title} {word}")
            queries.append(f"{slide_title} {word}")
    
    # Remove duplicates
    queries = list(set(queries))
    
    # If we don't have good queries, fall back to slide title
    if not queries or not any(q.strip() for q in queries):
        queries = [slide_title]
    
    return queries

def add_images_to_presentation(prs, outline, presentation_title, detail_level, num_images=3):
    """
    Add relevant images to the presentation based on the specified requirements
    """
    if not UNSPLASH_ACCESS_KEY:
        print("❌ Unsplash access key not found. Skipping image addition.")
        return prs
    
    # Create a temporary directory for downloaded images
    with tempfile.TemporaryDirectory() as temp_dir:
        if detail_level == "detailed":
            print("🔍 Processing detailed presentation with image slides...")
            
            # For detailed presentations: create image slides and insert them between content slides
            # Skip first 3 slides and last slide
            total_slides = len(prs.slides)
            
            # Determine where to insert image slides (after every 2-3 content slides, starting from slide 4)
            insert_after_indices = []
            for i in range(3, total_slides - 1, 3):  # Insert after every 3rd slide starting from slide 4
                if i < total_slides - 1:  # Don't insert after last slide
                    insert_after_indices.append(i)
            
            # Limit to 3-4 image slides max
            insert_after_indices = insert_after_indices[:4]
            
            print(f"📊 Will insert image slides after indices: {insert_after_indices}")
            
            # Process slides and create image slides
            image_slides_info = []  # Store (insert_after_index, image_path, title)
            
            # Create a list of outline sections with their indices
            outline_items = list(outline.items())
            
            for insert_after_index in insert_after_indices:
                if insert_after_index < len(outline_items):
                    section, content = outline_items[insert_after_index]
                    if content:
                        queries = get_relevant_image_queries(presentation_title, section, content, is_detailed=True)
                        
                        print(f"🔎 Searching images for slide '{section}' with queries: {queries}")
                        
                        for query in queries:
                            images = search_images(query, num_images)
                            if images:
                                # Select the most relevant image
                                selected_image = images[0]
                                image_path = os.path.join(temp_dir, f"{section}_{selected_image['id']}.jpg")
                                
                                if download_image(selected_image['download_url'], image_path):
                                    image_slides_info.append((insert_after_index, image_path, f"Visual: {section}"))
                                    print(f"✅ Found image for '{section}': {selected_image['description']}")
                                    break
            
            # Create image slides and insert them in the correct positions
            # We need to work backwards to maintain correct indices
            image_slides_info.sort(reverse=True, key=lambda x: x[0])  # Sort by index in descending order
            
            for insert_after_index, image_path, title in image_slides_info:
                # Create the image slide (this will be added at the end initially)
                image_slide = create_image_slide(prs, image_path, title)
                if image_slide:
                    print(f"🖼️ Created image slide for '{title}'")
                    
                    # Get the XML element of the newly created slide (last in the list)
                    new_slide_xml = prs.slides._sldIdLst[-1]
                    
                    # Remove it from the end
                    prs.slides._sldIdLst.remove(new_slide_xml)
                    
                    # Insert it at the correct position (after the specified slide)
                    insert_position = insert_after_index + 1
                    prs.slides._sldIdLst.insert(insert_position, new_slide_xml)
                    
                    print(f"📋 Inserted image slide at position {insert_position}")
                
        else:
            # For simple presentations: add smaller images to existing slides
            print("🔍 Processing simple presentation with inline images...")
            
            # Skip first 3 slides and last slide
            total_slides = len(outline)
            
            # Select slides for images (only slides 4 to n-1)
            eligible_slide_indices = list(range(3, total_slides - 1))
            
            # Select 3-4 slides for images based on presentation length
            num_slides_with_images = min(max(3, total_slides // 3), 4, len(eligible_slide_indices))
            
            if num_slides_with_images > 0:
                # Select slides for images (spread out)
                selected_indices = []
                if len(eligible_slide_indices) > num_slides_with_images:
                    step = len(eligible_slide_indices) // num_slides_with_images
                    for i in range(num_slides_with_images):
                        idx = min(i * step, len(eligible_slide_indices) - 1)
                        selected_indices.append(eligible_slide_indices[idx])
                else:
                    selected_indices = eligible_slide_indices
                
                print(f"📊 Adding images to slides at indices: {selected_indices}")
                
                # Process each selected slide
                for i, (section, content) in enumerate(outline.items()):
                    if i in selected_indices and content:
                        queries = get_relevant_image_queries(presentation_title, section, content, is_detailed=False)
                        
                        print(f"🔎 Searching images for slide '{section}' with queries: {queries}")
                        
                        for query in queries:
                            images = search_images(query, num_images)
                            if images:
                                # Select the most relevant image
                                selected_image = images[0]
                                image_path = os.path.join(temp_dir, f"{section}_{selected_image['id']}.jpg")
                                
                                if download_image(selected_image['download_url'], image_path):
                                    slide = prs.slides[i]
                                    add_image_to_slide(slide, image_path, "left_bottom", Inches(3))
                                    print(f"✅ Added image to slide '{section}': {selected_image['description']}")
                                    break
    
    return prs