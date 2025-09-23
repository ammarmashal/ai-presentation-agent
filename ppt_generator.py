import os
from dotenv import load_dotenv
from groq import Groq
import pprint
import re
import random
import argparse
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import tempfile
from collections import Counter
from llm_utils import generate_outline
import json
from pptx.enum.dml import MSO_LINE, MSO_FILL_TYPE, MSO_FILL
from PIL import Image
from PIL import Image as PILImage
import io

# Load environment variables
load_dotenv()
UNSPLASH_ACCESS_KEY = os.getenv("UNSPLASH_ACCESS_KEY")
GROQ_API_KEY = os.getenv("GROQ_API_KEY")




def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple"""
    try:
        # Remove # if present
        hex_color = hex_color.lstrip('#')
        # Convert to RGB
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    except:
        return (44, 62, 80)  # Default dark blue
    
    
def initialize_groq_client():
    """Initialize and return the Groq client"""
    if not GROQ_API_KEY:
        raise ValueError("❌ GROQ_API_KEY not found. Please add it to your .env file.")
    
    client = Groq(api_key=GROQ_API_KEY)
    return client

def is_valid_hex(s):
    if not s: return False
    s = s.strip()
    return bool(re.fullmatch(r'#?[0-9a-fA-F]{6}', s))

def normalize_hex(s):
    s = s.strip()
    if not s.startswith('#'):
        s = '#' + s
    return s.lower()

def resolve_color_value(val):
    """
    Accepts either a hex string like '#ff00aa' or a predefined color name
    from PREDEFINED_COLORS (case-insensitive, spaces allowed).
    Returns normalized hex string or None.
    """
    if not val:
        return None
    v = str(val).strip()
    # direct hex
    if is_valid_hex(v):
        return normalize_hex(v)
    # try mapping to PREDEFINED_COLORS (allow "Deep Space" or "deep_space")
    key = v.replace(' ', '_').lower()
    if key in PREDEFINED_COLORS:
        return PREDEFINED_COLORS[key]
    # maybe already lower-case key without underscore
    if v.lower() in PREDEFINED_COLORS:
        return PREDEFINED_COLORS[v.lower()]
    return None

# ----------------------------------------
# Improved calculate_text_color (WCAG-based)
# ----------------------------------------
def relative_luminance(rgb):
    def channel(c):
        c = c / 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    R, G, B = rgb
    return 0.2126 * channel(R) + 0.7152 * channel(G) + 0.0722 * channel(B)

def contrast_ratio(rgb1, rgb2):
    L1 = relative_luminance(rgb1)
    L2 = relative_luminance(rgb2)
    lighter = max(L1, L2)
    darker = min(L1, L2)
    return (lighter + 0.05) / (darker + 0.05)

def calculate_text_color(background_hex):
    """
    Return '#ffffff' or '#000000' depending on which gives better contrast.
    Prefer whichever meets WCAG 4.5:1 for normal text; if neither, pick the higher contrast one.
    """
    try:
        bg_rgb = hex_to_rgb(background_hex)
        white = (255, 255, 255)
        black = (0, 0, 0)
        contrast_white = contrast_ratio(bg_rgb, white)
        contrast_black = contrast_ratio(bg_rgb, black)
        # Prefer white if it meets 4.5 or has higher contrast
        if contrast_white >= 4.5 or contrast_white >= contrast_black:
            return "#ffffff"
        else:
            return "#000000"
    except Exception:
        return "#ffffff"
# Add this color dictionary near the top of the file, after imports
def relative_luminance(rgb):
    def channel(c):
        c = c / 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    R, G, B = rgb
    return 0.2126 * channel(R) + 0.7152 * channel(G) + 0.0722 * channel(B)

def contrast_ratio(rgb1, rgb2):
    L1 = relative_luminance(rgb1)
    L2 = relative_luminance(rgb2)
    lighter = max(L1, L2)
    darker = min(L1, L2)
    return (lighter + 0.05) / (darker + 0.05)

def calculate_text_color(background_hex):
    """
    Return '#ffffff' or '#000000' depending on which gives better contrast.
    Prefer whichever meets WCAG 4.5:1 for normal text; if neither, pick the higher contrast one.
    """
    try:
        bg_rgb = hex_to_rgb(background_hex)
        white = (255, 255, 255)
        black = (0, 0, 0)
        contrast_white = contrast_ratio(bg_rgb, white)
        contrast_black = contrast_ratio(bg_rgb, black)
        # Prefer white if it meets 4.5 or has higher contrast
        if contrast_white >= 4.5 or contrast_white >= contrast_black:
            return "#ffffff"
        else:
            return "#000000"
    except Exception:
        return "#ffffff"


# ----------------------------------------
# Robust LLM parsing without forcing 'dark'
# ----------------------------------------
def get_color_palette_from_llm(client, presentation_title, outline):
    """
    Use LLM to generate a color palette. Accept light or dark primary colors.
    """
    # Build outline text
    outline_text = f"Title: {presentation_title}\nSlides:\n"
    for section, points in outline.items():
        bullet_points = " • ".join([point["text"] for point in points[:3]])
        outline_text += f"- {section}: {bullet_points}\n"

    prompt = f"""
    You are an expert graphic designer. Based on the following presentation content, suggest a color palette.
    Analyze theme, mood, and key subjects.

    RETURN YOUR ANSWER *ONLY* AS A JSON OBJECT WITH THIS EXACT STRUCTURE - NO COMMENTS, NO EXPLANATIONS:
    {{
      "primary_color": "#xxxxxx" OR "color_name",
      "accent_color": "#xxxxxx" OR "color_name"
    }}

    The primary color can be dark or light; we will automatically set text color for readability.
    PRESENTATION CONTENT:
    {outline_text}
    """

    try:
        response = client.chat.completions.create(
            model="llama-3.1-8b-instant",
            messages=[
                {"role": "system", "content": "You are a helpful design assistant that suggests color palettes for presentations. Always respond with ONLY valid JSON, no comments or explanations."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            timeout=30
        )

        response_text = response.choices[0].message.content
        print(f"LLM Color Response: {response_text}")

        color_data = None
        # try direct json parse first
        try:
            color_data = json.loads(response_text)
        except Exception:
            # fallback: extract first JSON object in the text
            m = re.search(r'\{.*?\}', response_text, re.DOTALL)
            if m:
                try:
                    color_data = json.loads(m.group())
                except Exception as e:
                    print(f"❌ JSON parse error after regex: {e}")
                    color_data = None

        # final fallback or validation
        if not color_data or not isinstance(color_data, dict):
            print("⚠️ No valid JSON palette found in LLM response. Using defaults.")
            color_data = {}

        # Resolve values (hex or predefined name)
        primary_hex = resolve_color_value(color_data.get("primary_color"))
        accent_hex = resolve_color_value(color_data.get("accent_color"))

        # Fallback defaults only if unresolved
        if not primary_hex:
            primary_hex = "#2c3e50"  # default primary
        if not accent_hex:
            accent_hex = "#3498db"  # default accent

        print(f"🎨 Selected colors: Primary: {primary_hex}, Accent: {accent_hex}")
        return {"primary_color": primary_hex, "accent_color": accent_hex}

    except Exception as e:
        print(f"❌ Error getting color palette from LLM: {e}")
        return {"primary_color": "#2c3e50", "accent_color": "#3498db"}

# ----------------------------------------
# apply_color_theme stays mostly same
# ----------------------------------------
def apply_color_theme(prs, color_palette):
    """Apply color theme to presentation background (primary_color expected as hex)"""
    primary_rgb = hex_to_rgb(color_palette["primary_color"])

    slide_master = prs.slide_masters[0]
    for slide_layout in slide_master.slide_layouts:
        try:
            background = slide_layout.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*primary_rgb)
        except Exception:
            continue

    return primary_rgb


def search_images(query, count=5):
    """
    Search for images using Unsplash API
    """
    if not UNSPLASH_ACCESS_KEY:
        print("❌ Unsplash access key not found. Please add UNSPLASH_ACCESS_KEY to your .env file.")
        return []
    
    try:
        url = f"https://api.unsplash.com/search/photos"
        params = {
            "query": query,
            "per_page": count,
            "client_id": UNSPLASH_ACCESS_KEY
        }
        
        response = requests.get(url, params=params, timeout=30)
        response.raise_for_status()
        
        data = response.json()
        images = []
        
        for result in data.get("results", []):
            images.append({
                "id": result["id"],
                "url": result["urls"]["regular"],
                "description": result.get("description", result.get("alt_description", "")),
                "download_url": result["urls"]["regular"]
            })
        
        return images
    
    except Exception as e:
        print(f"❌ Error searching images: {e}")
        return []

def download_image(image_url, filename):
    """
    Download an image from a URL
    """
    try:
        response = requests.get(image_url, timeout=30)
        response.raise_for_status()
        
        with open(filename, 'wb') as f:
            f.write(response.content)
        
        return True
    except Exception as e:
        print(f"❌ Error downloading image: {e}")
        return False

def get_relevant_image_queries(title, section, points, is_detailed):
    queries = []
    queries.append(title)
    queries.append(f"{section} {title}")
    extra_keywords = ["illustration", "concept", "background", "design", "education"]
    for kw in extra_keywords:
        queries.append(f"{title} {section} {kw}")
    if points:
        first_point = points[0]["text"] if isinstance(points[0], dict) else str(points[0])
        keywords = " ".join(first_point.split()[:3])
        queries.append(f"{title} {keywords}")
    return queries


def determine_slide_layout(slide_index, detail_level, content_length, outline, section_title):
    """
    Determine the appropriate layout for each slide based on its position, content, and title.
    """
    num_slides = len(outline)

    # 1. Prioritize the Title Slide
    if slide_index == 0:
        return "title"

    # 2. Prioritize the Conclusion Slide
    if "conclusion" in section_title.lower():
        return "conclusion"
    
    # 3. Handle the Introduction Slide
    if slide_index == 1:
        return "title_content"

    # Available layouts, excluding 'image_full' and 'comparison'
    image_layouts = ["image_left_text_right", "image_right_text_left","title"]
    content_layouts = ["title_content", "two_column"]
    
    # Combine lists to create a full set of options
    all_layouts = image_layouts + content_layouts
    
    # Apply logic based on detail level and content length
    if detail_level == "detailed":
        # For detailed, prefer layouts with more space for text
        layout_options = content_layouts + image_layouts
    else:
        # For simple, vary the layouts more evenly
        layout_options = all_layouts

    # Use more image layouts for slides with less content
    if content_length <= 2 and image_layouts:
        layout_options = image_layouts + content_layouts # Prioritize image layouts
    
    # Use more content-focused layouts for slides with more content
    if content_length >= 5 and content_layouts:
        layout_options = content_layouts + image_layouts # Prioritize content layouts
        
    return random.choice(layout_options)


def create_image_title_slide(prs, presentation_title, image_path, primary_rgb, text_rgb, accent_rgb):
    """
    Creates a title slide with a full-page background image and a text box
    in the bottom-right corner with a fixed size.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # --- Add Image as a Full-page Background ---
    img_left = Inches(0)
    img_top = Inches(0)
    img_width = slide_width
    img_height = slide_height

    try:
        with PILImage.open(image_path) as img:
            original_width_px, original_height_px = img.size
        
        container_aspect_ratio = float(img_width.emu) / float(img_height.emu)
        image_aspect_ratio = float(original_width_px) / float(original_height_px)

        if image_aspect_ratio > container_aspect_ratio:
            new_height_emu = img_height.emu
            new_width_emu = int(new_height_emu * image_aspect_ratio)
            picture_left_emu = img_left.emu + (img_width.emu - new_width_emu) / 2
            picture_top_emu = img_top.emu
        else:
            new_width_emu = img_width.emu
            new_height_emu = int(new_width_emu / image_aspect_ratio)
            picture_left_emu = img_left.emu
            picture_top_emu = img_top.emu + (img_height.emu - new_height_emu) / 2

        slide.shapes.add_picture(
            image_path, picture_left_emu, picture_top_emu,
            width=new_width_emu, height=new_height_emu
        )
    
    except FileNotFoundError:
        print(f"Error: Image file not found at {image_path}. Using solid background.")
        background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height)
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(*primary_rgb)
        background.line.fill.background()

    # --- Add Text in a Bounding Box in the Bottom-Right Corner ---
    # CHANGED: Using slightly smaller fixed dimensions
    box_width = Inches(4.0)
    box_height = Inches(1.8)
    box_left = slide_width - box_width - Inches(0.3)
    box_top = slide_height - box_height - Inches(0.2)

    # Create the white bounding box
    text_container = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, box_left, box_top, box_width, box_height
    )
    text_container.fill.solid()
    text_container.fill.fore_color.rgb = RGBColor(255, 255, 255)
    text_container.line.fill.background()

    # Add text box inside the white container
    # The text box is slightly smaller than the container to provide padding
    text_box = slide.shapes.add_textbox(box_left + Inches(0.2), box_top + Inches(0.2), 
                                         box_width - Inches(0.4), box_height - Inches(0.4))
    title_frame = text_box.text_frame
    title_frame.word_wrap = True
    title_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    title_frame.clear()
    
    p_title = title_frame.paragraphs[0]
    p_title.text = presentation_title
    p_title.alignment = PP_ALIGN.CENTER
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(0, 0, 0)

    p_subtitle = title_frame.add_paragraph()
    p_subtitle.text = "Generated with AI Presentation Generator"
    p_subtitle.alignment = PP_ALIGN.CENTER
    p_subtitle.font.size = Pt(14)
    p_subtitle.font.color.rgb = RGBColor(0, 0, 0)
    p_subtitle.space_before = Inches(0.3)

    return slide


'''def create_image_title_slide(prs, presentation_title, image_path, primary_rgb, text_rgb, accent_rgb):
    """
    Creates a title slide with an image on the left and title text on the right,
    mimicking a modern, stylish layout.
    """
    # Use a blank slide layout (index 6) for a custom layout
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Get slide dimensions for proportional scaling
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # --- Set Background for the Right Half (Gradient) ---
    right_half_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        slide_width / 2, 
        Inches(0), 
        slide_width / 2, 
        slide_height
    )
    fill = right_half_shape.fill

    # Initialize gradient fill
    fill.gradient()
    
    # Set gradient properties
    fill.gradient_angle = 180

    # Handle gradient stops
    start_color_gradient = RGBColor(255, 255, 255)  # White at the top
    end_color_gradient = RGBColor(*primary_rgb)
    
    # Ensure we have at least 2 gradient stops
    while len(fill.gradient_stops) < 2:
        fill.gradient_stops.add_stop(0.0, RGBColor(0, 0, 0))
    
    # Modify the first two stops
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[0].color.rgb = start_color_gradient
    
    fill.gradient_stops[1].position = 1.0
    fill.gradient_stops[1].color.rgb = end_color_gradient
    
    # Remove any additional stops if they exist
    for i in range(len(fill.gradient_stops) - 1, 1, -1):
        fill.gradient_stops[i].delete()
    
    right_half_shape.line.fill.background()

    # --- Add Image on the Left ---
    img_left = Inches(0)
    img_top = Inches(0)
    img_width = slide_width / 2
    img_height = slide_height

    try:
        with PILImage.open(image_path) as img:
            original_width_px, original_height_px = img.size
        
        container_aspect_ratio = float(img_width) / float(img_height)
        image_aspect_ratio = float(original_width_px) / float(original_height_px)

        if image_aspect_ratio > container_aspect_ratio:
            new_height_emu = img_height
            new_width_emu = int(img_height * image_aspect_ratio)
            picture_left_emu = img_left.emu + (img_width - new_width_emu) / 2
            picture_top_emu = img_top.emu
        else:
            new_width_emu = img_width
            new_height_emu = int(img_width / image_aspect_ratio)
            picture_left_emu = img_left.emu
            picture_top_emu = img_top.emu + (img_height - new_height_emu) / 2

        slide.shapes.add_picture(
            image_path, picture_left_emu, picture_top_emu,
            width=new_width_emu, height=new_height_emu
        )
    
    except FileNotFoundError:
        print(f"Error: Image file not found at {image_path}. Using a grey placeholder.")
        placeholder_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, img_left, img_top, img_width, img_height
        )
        placeholder_fill = placeholder_shape.fill
        placeholder_fill.solid()
        placeholder_fill.fore_color.rgb = RGBColor(180, 180, 180)
        placeholder_shape.line.fill.background()
    
    # --- Add Title and Subtitle Text on the Right ---
    title_box_left = slide_width / 2 + Inches(0.5)
    title_box_top = Inches(1.5)
    title_box_width = slide_width / 2 - Inches(1)
    title_box_height = Inches(3)

    title_box = slide.shapes.add_textbox(title_box_left, title_box_top, title_box_width, title_box_height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    
    # CORRECTED: Simple and reliable way to clear and add text
    # Clear all text first
    title_frame.clear()
    
    # Add title paragraph
    p_title = title_frame.paragraphs[0]
    p_title.text = presentation_title
    p_title.alignment = PP_ALIGN.CENTER
    p_title.font.size = Pt(44)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(*accent_rgb)

    # Add subtitle paragraph
    p_subtitle = title_frame.add_paragraph()
    p_subtitle.text = "Generated with AI Presentation Generator'Ammar Yasser'"
    p_subtitle.alignment = PP_ALIGN.CENTER
    p_subtitle.font.size = Pt(18)
    p_subtitle.font.color.rgb = RGBColor(*text_rgb)
    p_subtitle.space_before = Inches(0.3)

    title_box.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return slide'''
# And a simple fallback create_title_slide (if not already defined and used in your main loop)
def create_title_slide(prs, presentation_title, text_rgb):
    """
    Creates a basic text-only title slide (fallback for image issues).
    """
    slide = prs.slides.add_slide(prs.slide_layouts[0]) # Use default title layout
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = presentation_title
    subtitle.text = "Created with AI Presentation Generator"
    
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.color.rgb = RGBColor(*text_rgb)
        paragraph.font.bold = True
        paragraph.font.size = Pt(44)
    
    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.font.color.rgb = RGBColor(*text_rgb)
        paragraph.font.size = Pt(18)
    
    return slide

def create_title_content_slide(prs, section, points, text_rgb):
    """Create a slide with title and content with theme colors"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = section
    
    # Apply text color to title
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.font.color.rgb = RGBColor(*text_rgb)
    
    text_frame = content_shape.text_frame
    text_frame.clear()
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.word_wrap = True
    
    first = True
    for point in points:
        level = point.get("level", 0)
        
        if first:
            p = text_frame.paragraphs[0]
            first = False
        else:
            p = text_frame.add_paragraph()
        
        p.text = point["text"]
        p.level = level
        p.font.color.rgb = RGBColor(*text_rgb)
    
    return slide

def create_image_left_text_right_slide(prs, section, points, image_path, primary_rgb, text_rgb, accent_rgb):
    """
    Creates a slide with image on left and text on right, with modern styling
    and dynamic theme colors.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Define layout dimensions
    text_left = Inches(5.2)
    text_top = Inches(1.5)
    text_width = Inches(4.3)
    text_height = Inches(4.5)
    img_left = Inches(0.8)
    img_top = Inches(1.5)
    img_width = Inches(4)
    img_height = Inches(4.5)

    # Add text container (the border)
    text_container = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, text_left, text_top, text_width, text_height
    )
    text_container.fill.solid() # CHANGED: Solid white background for the text container
    text_container.fill.fore_color.rgb = RGBColor(255, 255, 255)
    text_container.line.width = Pt(1)
    text_container.line.color.rgb = RGBColor(*accent_rgb) # Border line color is accent
    text_container.line.dash_style = MSO_LINE.DASH

    # Add title with dynamic color
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = section
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(*text_rgb)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add modern decorative line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.1), Inches(8), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(*accent_rgb)
    line.line.fill.background()

    # Add image container
    img_container = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, img_left, img_top, img_width, img_height
    )
    img_container.fill.solid()
    img_container.fill.fore_color.rgb = RGBColor(255, 255, 255)
    img_container.line.width = Pt(2)
    img_container.line.color.rgb = RGBColor(*accent_rgb)
    
    try:
        with Image.open(image_path) as img:
            original_width_px, original_height_px = img.size
    except FileNotFoundError:
        print(f"Error: Image file not found at {image_path}")
        placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, img_left, img_top, img_width, img_height)
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(200, 200, 200)
        placeholder.line.width = Pt(2)
        placeholder.line.color.rgb = RGBColor(*text_rgb)
    else:
        container_width_emu = img_width.emu
        container_height_emu = img_height.emu
        ratio_w = container_width_emu / original_width_px
        ratio_h = container_height_emu / original_height_px
        scale_factor = min(ratio_w, ratio_h)
        new_width_emu = int(original_width_px * scale_factor)
        new_height_emu = int(original_height_px * scale_factor)
        picture_left_emu = img_left.emu + (container_width_emu - new_width_emu) / 2
        picture_top_emu = img_top.emu + (container_height_emu - new_height_emu) / 2
        slide.shapes.add_picture(
            image_path, picture_left_emu, picture_top_emu, 
            width=new_width_emu, height=new_height_emu
        )

    # Add text box with modern styling
    text_inner_left = text_left + Inches(0.2)
    text_inner_top = text_top + Inches(0.5)
    text_inner_width = text_width - Inches(0.4)
    text_inner_height = text_height - Inches(0.4)
    text_box = slide.shapes.add_textbox(
        text_inner_left, text_inner_top, text_inner_width, text_inner_height
    )
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    first = True
    for i, point in enumerate(points):
        level = point.get("level", 0)
        p = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
        first = False
        
        p.text = "●  " + point["text"]
        p.level = level
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_after = Inches(0.15) if i < len(points) - 1 else Inches(0.1)

        if i == 0:
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)

    # Add decorative elements
    decor_right = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(0.1), Inches(0.4), Inches(0.4))
    decor_right.fill.solid()
    decor_right.fill.fore_color.rgb = RGBColor(*accent_rgb)
    decor_right.line.fill.background()
    
    decor_left = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.1), Inches(6.5), Inches(0.3), Inches(0.3))
    decor_left.fill.solid()
    decor_left.fill.fore_color.rgb = RGBColor(*primary_rgb)
    decor_left.line.fill.background()
    
    return slide



def create_image_right_text_left_slide(prs, section, points, image_path, primary_rgb, text_rgb, accent_rgb):
    """
    Creates a slide with image on right and text on left, with modern styling
    and dynamic theme colors.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Define layout dimensions
    text_left = Inches(0.5)
    text_top = Inches(1.5)
    text_width = Inches(4.5)
    text_height = Inches(4.5)
    img_left = Inches(5.2)
    img_top = Inches(1.5)
    img_width = Inches(4.3)
    img_height = Inches(4.5)

    # Add title with dynamic color
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = section
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(*text_rgb)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add modern decorative line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.1), Inches(8), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(*accent_rgb)
    line.line.fill.background()

    # Add text container (the border)
    text_container = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, text_left, text_top, text_width, text_height
    )
    text_container.fill.solid() # CHANGED: Solid white background for the text container
    text_container.fill.fore_color.rgb = RGBColor(255, 255, 255)
    text_container.line.width = Pt(1)
    text_container.line.color.rgb = RGBColor(*accent_rgb) # Border line color is accent
    text_container.line.dash_style = MSO_LINE.DASH

    # Add text box with bullet points, positioned inside the container
    text_inner_left = text_left + Inches(0.2)
    text_inner_top = text_top + Inches(0.5)
    text_inner_width = text_width - Inches(0.4)
    text_inner_height = text_height - Inches(0.4)
    
    text_box = slide.shapes.add_textbox(text_inner_left, text_inner_top, text_inner_width, text_inner_height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    first = True
    for i, point in enumerate(points):
        level = point.get("level", 0)
        p = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
        first = False
        
        p.text = "●  " + point["text"]
        p.level = level
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0, 0, 0) # Fixed black color for body text
        p.space_after = Inches(0.15) if i < len(points) - 1 else Inches(0.1)

        if i == 0:
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0) # First point also has black color
            
    # Add image container
    img_container = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, img_left, img_top, img_width, img_height
    )
    img_container.fill.solid()
    img_container.fill.fore_color.rgb = RGBColor(255, 255, 255)
    img_container.line.width = Pt(2)
    img_container.line.color.rgb = RGBColor(*accent_rgb)
    
    try:
        with Image.open(image_path) as img:
            original_width_px, original_height_px = img.size
    except FileNotFoundError:
        print(f"Error: Image file not found at {image_path}")
        placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, img_left, img_top, img_width, img_height)
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(200, 200, 200)
        placeholder.line.width = Pt(2)
        placeholder.line.color.rgb = RGBColor(*text_rgb)
    else:
        container_width_emu = img_width.emu
        container_height_emu = img_height.emu
        ratio_w = container_width_emu / original_width_px
        ratio_h = container_height_emu / original_height_px
        scale_factor = min(ratio_w, ratio_h)
        new_width_emu = int(original_width_px * scale_factor)
        new_height_emu = int(original_height_px * scale_factor)
        picture_left_emu = img_left.emu + (container_width_emu - new_width_emu) / 2
        picture_top_emu = img_top.emu + (container_height_emu - new_height_emu) / 2
        slide.shapes.add_picture(
            image_path, picture_left_emu, picture_top_emu, 
            width=new_width_emu, height=new_height_emu
        )
    
    # Add decorative elements
    decor_right = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(0.1), Inches(0.4), Inches(0.4))
    decor_right.fill.solid()
    decor_right.fill.fore_color.rgb = RGBColor(*accent_rgb)
    decor_right.line.fill.background()
    
    decor_left = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.1), Inches(6.5), Inches(0.3), Inches(0.3))
    decor_left.fill.solid()
    decor_left.fill.fore_color.rgb = RGBColor(*primary_rgb)
    decor_left.line.fill.background()
    
    return slide

def create_two_column_slide(prs, section, points, text_rgb):
    """
    Creates a two-column text slide with modern styling, including borders,
    bullet points, and decorative elements.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # --- Define Layout Dimensions ---
    left_text_left = Inches(0.5)
    left_text_top = Inches(1.5)
    left_text_width = Inches(4.5)
    left_text_height = Inches(5.5)
    
    right_text_left = Inches(5.2)
    right_text_top = Inches(1.5)
    right_text_width = Inches(4.3)
    right_text_height = Inches(5.5)

    # --- Add Title and Decorative Line ---
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = section
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(*text_rgb)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.1), Inches(8), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(*text_rgb)
    line.line.fill.background()
    
    # --- Split points into two columns ---
    mid_point = len(points) // 2
    left_points = points[:mid_point]
    right_points = points[mid_point:]
    
    # --- Add Left Text Column Container and Content ---
    left_text_container = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left_text_left, left_text_top, left_text_width, left_text_height
    )
    left_text_container.fill.background()
    left_text_container.line.width = Pt(1)
    left_text_container.line.color.rgb = RGBColor(*text_rgb)
    left_text_container.line.dash_style = MSO_LINE.DASH
    
    left_text_box = slide.shapes.add_textbox(left_text_left + Inches(0.2), left_text_top + Inches(0.5), 
                                             left_text_width - Inches(0.4), left_text_height - Inches(0.4))
    left_text_frame = left_text_box.text_frame
    left_text_frame.word_wrap = True
    
    first = True
    for i, point in enumerate(left_points):
        p = left_text_frame.paragraphs[0] if first else left_text_frame.add_paragraph()
        first = False
        p.text = "•  " + point["text"]
        p.level = point.get("level", 0)
        p.font.size = Pt(20)  # Increased font size
        p.font.color.rgb = RGBColor(*text_rgb)
        p.space_after = Inches(0.15) if i < len(left_points) - 1 else Inches(0.1)
    
    # --- Add Right Text Column Container and Content ---
    right_text_container = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, right_text_left, right_text_top, right_text_width, right_text_height
    )
    right_text_container.fill.background()
    right_text_container.line.width = Pt(1)
    right_text_container.line.color.rgb = RGBColor(*text_rgb)
    right_text_container.line.dash_style = MSO_LINE.DASH

    right_text_box = slide.shapes.add_textbox(right_text_left + Inches(0.2), right_text_top + Inches(0.5), 
                                              right_text_width - Inches(0.4), right_text_height - Inches(0.4))
    right_text_frame = right_text_box.text_frame
    right_text_frame.word_wrap = True
    
    first = True
    for i, point in enumerate(right_points):
        p = right_text_frame.paragraphs[0] if first else right_text_frame.add_paragraph()
        first = False
        p.text = "•  " + point["text"]
        p.level = point.get("level", 0)
        p.font.size = Pt(20)  # Increased font size
        p.font.color.rgb = RGBColor(*text_rgb)
        p.space_after = Inches(0.15) if i < len(right_points) - 1 else Inches(0.1)

    # --- Add Decorative Elements ---
    decor_right = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(0.1), Inches(0.4), Inches(0.4))
    decor_right.fill.solid()
    decor_right.fill.fore_color.rgb = RGBColor(*text_rgb)
    decor_right.line.fill.background()
    
    decor_left = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.1), Inches(6.5), Inches(0.3), Inches(0.3))
    decor_left.fill.solid()
    decor_left.fill.fore_color.rgb = RGBColor(*text_rgb)
    decor_left.line.fill.background()
    
    return slide

def create_conclusion_slide(prs, section, points, text_rgb):
    """Create a conclusion slide with theme colors"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = section
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.color.rgb = RGBColor(*text_rgb)
    
    # Add content in the center
    content_left = Inches(1)
    content_top = Inches(2)
    content_width = Inches(8)
    content_height = Inches(4)
    content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Add key points
    for i, point in enumerate(points):
        if i < 3:  # Limit to 3 key points
            p = content_frame.add_paragraph() if i > 0 else content_frame.paragraphs[0]
            p.text = point["text"]
            p.font.size = Pt(20)
            p.level = 0
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = RGBColor(*text_rgb)
    
    return slide

def create_thank_you_slide(prs, primary_rgb, text_rgb, accent_rgb):
    """Create a Thank You slide with theme colors"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add background color using theme's primary color
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(*primary_rgb)
    background.line.fill.background()
    
    # Add Thank You text
    thank_you_left = Inches(1)
    thank_you_top = Inches(2.5)
    thank_you_width = Inches(8)
    thank_you_height = Inches(2)
    thank_you_box = slide.shapes.add_textbox(thank_you_left, thank_you_top, thank_you_width, thank_you_height)
    thank_you_frame = thank_you_box.text_frame
    thank_you_frame.text = "Thank You!"
    thank_you_frame.paragraphs[0].font.size = Pt(48)
    thank_you_frame.paragraphs[0].font.bold = True
    thank_you_frame.paragraphs[0].font.color.rgb = RGBColor(*text_rgb)
    thank_you_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add smaller subtitle
    subtitle_left = Inches(1)
    subtitle_top = Inches(5)
    subtitle_width = Inches(8)
    subtitle_height = Inches(1)
    subtitle_box = slide.shapes.add_textbox(subtitle_left, subtitle_top, subtitle_width, subtitle_height)
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Questions & Discussion"
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    # Use accent color for subtitle or lighter version of text color
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(*accent_rgb)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def list_available_themes(theme_folder="themes"):
    """List all available themes in the specified folder"""
    if not os.path.exists(theme_folder):
        os.makedirs(theme_folder)
        return []
    
    themes = [f for f in os.listdir(theme_folder) if f.endswith(".pptx")]
    return themes

def get_theme_path(theme_name, theme_folder="themes"):
    """Get the full path to a theme file"""
    return os.path.join(theme_folder, f"{theme_name}.pptx")

def create_presentation(outline, presentation_title, detail_level, filename="presentation.pptx", theme_path=None):
    """
    Creates a PowerPoint presentation from a structured outline with dynamic color theming.
    """
    try:
        # Initialize Groq client and get color palette
        client = initialize_groq_client()
        color_palette = get_color_palette_from_llm(client, presentation_title, outline)
        primary_rgb = hex_to_rgb(color_palette["primary_color"])
        accent_rgb = hex_to_rgb(color_palette["accent_color"])
        text_hex = calculate_text_color(color_palette["primary_color"])
        text_rgb = hex_to_rgb(text_hex)
        prs = Presentation()
        apply_color_theme(prs, color_palette)
        
        # Ensure the presentation ends with a conclusion slide
        image_folder = "Img"
        if not os.path.exists(image_folder):
            os.makedirs(image_folder)
            print(f"✅ Created image folder: {image_folder}")
        
        # --- START: Logic for title slide with an image ---
        print("Creating title slide with image...")
        
        title_queries = get_relevant_image_queries(
            presentation_title, 
            "Introduction", 
            list(outline.values())[0], 
            detail_level == "detailed"
        )
        title_image_path = None
        
        for query in title_queries:
            images = search_images(query, 5)  # CHANGED: Search for more images
            if images:
                # CHANGED: Pick a random image from the results
                image_to_download = random.choice(images)
                
                sanitized_title = re.sub(r'[\\/:*?"<>|]', '_', presentation_title)
                title_image_path = os.path.join(image_folder, f"{sanitized_title}_{image_to_download['id']}.jpg")
                if download_image(image_to_download['download_url'], title_image_path):
                    break
                else:
                    title_image_path = None
        
        if title_image_path and os.path.exists(title_image_path):
            create_image_title_slide(
                prs, 
                presentation_title, 
                title_image_path, 
                primary_rgb,
                text_rgb,
                accent_rgb
            )
        else:
            print(f"⚠️ No image found for title slide, using text-only layout.")
            create_title_slide(prs, presentation_title, text_rgb)
        
        # --- END: Logic for title slide with an image ---

        # Iterate through the outline to create each slide
        for i, (section, points) in enumerate(outline.items()):
            slide_index = i + 1
            layout = determine_slide_layout(slide_index, detail_level, len(points), outline, section)
            
            print(f"Creating slide {slide_index}: {section} ({layout} layout)")
            
            slide = None
            if layout == "title_content":
                slide = create_title_content_slide(prs, section, points, text_rgb)
            
            elif layout in ["image_left_text_right", "image_right_text_left", "image_full"]:
                queries = get_relevant_image_queries(presentation_title, section, points, detail_level == "detailed")
                image_found = False
                
                for query in queries:
                    images = search_images(query, 5) # CHANGED: Search for more images
                    if images:
                        # CHANGED: Pick a random image from the results
                        image_to_download = random.choice(images)
                        
                        sanitized_section = re.sub(r'[\\/:*?"<>|]', '_', section)
                        image_path = os.path.join(image_folder, f"{sanitized_section}_{image_to_download['id']}.jpg")
                        if download_image(image_to_download['download_url'], image_path):
                            image_found = True
                            if layout == "image_left_text_right":
                                slide = create_image_left_text_right_slide(prs, section, points, image_path, primary_rgb, text_rgb, accent_rgb)
                            elif layout == "image_right_text_left":
                                slide = create_image_right_text_left_slide(prs, section, points, image_path, primary_rgb, text_rgb, accent_rgb)
                            break
                
                if not image_found:
                    print(f"⚠️ No image found for '{section}', using title_content layout instead.")
                    slide = create_title_content_slide(prs, section, points, text_rgb)
            
            elif layout == "two_column":
                slide = create_two_column_slide(prs, section, points, text_rgb)
            
            elif layout == "conclusion":
                slide = create_conclusion_slide(prs, section, points, text_rgb)
            
            # Apply background color explicitly as fallback
            if slide:
                try:
                    background = slide.background
                    fill = background.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(*primary_rgb)
                except Exception:
                    pass
        
        # Add Thank You slide at the end with theme colors
        create_thank_you_slide(prs, primary_rgb, text_rgb, accent_rgb)
        
        # Save the completed presentation
        prs.save(filename)
        print(f"✅ Presentation saved as: {filename}")
        return filename
    
    except Exception as e:
        print(f"❌ Error creating presentation: {e}")
        raise


def ensure_conclusion_slide(outline, presentation_title):
    """Ensure the presentation always ends with a conclusion slide"""
    last_section = list(outline.keys())[-1] if outline else ""
    if not any(keyword in last_section.lower() for keyword in ["conclusion", "summary", "wrap-up", "final"]):
        outline["Conclusion"] = [
            {"text": "Summary of key points and takeaways", "level": 0},
            {"text": "Future outlook and recommendations", "level": 0},
            {"text": "Q&A and discussion", "level": 0}
        ]
    return outline


def main():
    parser = argparse.ArgumentParser(description="Generate AI-powered presentations with dynamic color theming")
    parser.add_argument("--topic", type=str, default="Artificial Intelligence", help="Presentation topic")
    parser.add_argument("--detail", type=str, default="simple", choices=["simple", "detailed"], help="Detail level")
    parser.add_argument("--theme", type=str, help="Name of the theme file (without extension)")
    parser.add_argument("--output", type=str, default="presentation.pptx", help="Output filename")
    
    args = parser.parse_args()
    
    theme_folder = "themes"
    available_themes = list_available_themes(theme_folder)
    
    print("🎯 AI Presentation Generator with Dynamic Color Theming")
    print("=" * 50)
    print(f"Topic: {args.topic}")
    print(f"Detail Level: {args.detail}")
    
    if available_themes:
        print("Available Themes:")
        for theme in available_themes:
            print(f" - {os.path.splitext(theme)[0]}")
        
        if args.theme:
            theme_path = get_theme_path(args.theme, theme_folder)
            if not os.path.exists(theme_path):
                print(f"❌ Theme '{args.theme}' not found. Using dynamic color theming.")
                theme_path = None
        else:
            print("ℹ️  No theme specified. Using dynamic color theming.")
            theme_path = None
    else:
        print("ℹ️  No themes found. Using dynamic color theming.")
        theme_path = None
    
    print("=" * 50)
    
    outline, presentation_title = generate_outline(args.topic, args.detail)
    outline = ensure_conclusion_slide(outline, presentation_title)
    create_presentation(outline, presentation_title, args.detail, args.output, theme_path)
    print("🎉 Presentation created successfully with dynamic color theming!")


if __name__ == "__main__":
    main()
