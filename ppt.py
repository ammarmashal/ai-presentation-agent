import os
from dotenv import load_dotenv
from groq import Groq
import pprint
import re
import random
import argparse
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import tempfile
from collections import Counter
from llm_utils import generate_outline



# Load environment variables
load_dotenv()
UNSPLASH_ACCESS_KEY = os.getenv("UNSPLASH_ACCESS_KEY")


def search_images(query, count=5):
    """
    Search for images using Unsplash API
    """
    if not UNSPLASH_ACCESS_KEY:
        print("❌ Unsplash access key not found. Please add UNSPLASH_ACCESS_KEY to your .env file.")
        return []
    
    try:
        url = f"https://api.unsplash.com/search/photos"
        params = {
            "query": query,
            "per_page": count,
            "client_id": UNSPLASH_ACCESS_KEY
        }
        
        response = requests.get(url, params=params, timeout=30)
        response.raise_for_status()
        
        data = response.json()
        images = []
        
        for result in data.get("results", []):
            images.append({
                "id": result["id"],
                "url": result["urls"]["regular"],
                "description": result.get("description", result.get("alt_description", "")),
                "download_url": result["urls"]["regular"]
            })
        
        return images
    
    except Exception as e:
        print(f"❌ Error searching images: {e}")
        return []

def download_image(image_url, filename):
    """
    Download an image from a URL
    """
    try:
        response = requests.get(image_url, timeout=30)
        response.raise_for_status()
        
        with open(filename, 'wb') as f:
            f.write(response.content)
        
        return True
    except Exception as e:
        print(f"❌ Error downloading image: {e}")
        return False

def get_relevant_image_queries(presentation_title, slide_title, content, is_detailed=False):
    """
    Generate relevant search queries for images based on presentation content
    For detailed presentations, use the entire LLM content
    """
    queries = []
    
    if is_detailed:
        # For detailed presentations, use the entire content from LLM
        content_text = " ".join([point["text"] for point in content])
        all_text = f"{presentation_title} {slide_title} {content_text}"
        
        # Extract meaningful words (nouns, adjectives)
        words = re.findall(r'\b[a-zA-Z]{4,}\b', all_text.lower())
        
        # Remove common stop words
        stop_words = {"the", "and", "or", "but", "in", "on", "at", "to", "for", 
                    "of", "with", "by", "that", "this", "these", "those", "is",
                    "are", "was", "were", "be", "been", "being", "have", "has",
                    "had", "do", "does", "did", "will", "would", "could", "should"}
        
        meaningful_words = [word for word in words if word not in stop_words]
        
        # Get the most frequent words
        word_counts = Counter(meaningful_words)
        most_common = [word for word, count in word_counts.most_common(8)]
        
        queries.extend(most_common)
        
        # Add combinations
        for word in most_common[:3]:
            queries.append(f"{presentation_title} {word}")
            queries.append(f"{slide_title} {word}")
    else:
        # For simple presentations, use a more targeted approach
        content_text = " ".join([point["text"] for point in content])
        all_text = f"{slide_title} {content_text}"
        
        # Find specific nouns (more likely to have good images)
        words = re.findall(r'\b[a-zA-Z]{4,}\b', all_text.lower())
        
        # Filter out abstract terms
        abstract_terms = {"management", "strategy", "system", "process", "method", 
                         "approach", "concept", "theory", "principle", "framework",
                         "model", "analysis", "development", "implementation"}
        
        concrete_words = [word for word in words if word not in abstract_terms]
        
        # Add the most frequent concrete words
        word_counts = Counter(concrete_words)
        most_common = [word for word, count in word_counts.most_common(5)]
        
        queries.extend(most_common)
        
        # Add combinations with presentation title for context
        for word in most_common[:2]:
            queries.append(f"{presentation_title} {word}")
            queries.append(f"{slide_title} {word}")
    
    # Remove duplicates
    queries = list(set(queries))
    
    # If we don't have good queries, fall back to slide title
    if not queries or not any(q.strip() for q in queries):
        queries = [slide_title]
    
    return queries

def determine_slide_layout(slide_index, detail_level, content_length, outline, section_title):
    """
    Determine the appropriate layout for each slide based on its position, content, and title.
    """
    num_slides = len(outline)

    # 1. Prioritize the Title Slide
    if slide_index == 0:
        return "title"

    # 2. Prioritize the Conclusion Slide
    if "conclusion" in section_title.lower():
        return "conclusion"
    
    # 3. Handle the Introduction Slide
    if slide_index == 1:
        return "title_content"

    # Available layouts, excluding 'image_full' and 'comparison'
    image_layouts = ["image_left_text_right", "image_right_text_left"]
    content_layouts = ["title_content", "two_column"]
    
    # Combine lists to create a full set of options
    all_layouts = image_layouts + content_layouts
    
    # Apply logic based on detail level and content length
    if detail_level == "detailed":
        # For detailed, prefer layouts with more space for text
        layout_options = content_layouts + image_layouts
    else:
        # For simple, vary the layouts more evenly
        layout_options = all_layouts

    # Use more image layouts for slides with less content
    if content_length <= 2 and image_layouts:
        layout_options = image_layouts + content_layouts # Prioritize image layouts
    
    # Use more content-focused layouts for slides with more content
    if content_length >= 5 and content_layouts:
        layout_options = content_layouts + image_layouts # Prioritize content layouts
        
    return random.choice(layout_options)

def create_title_slide(prs, presentation_title):
    """Create a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = presentation_title
    subtitle.text = "Created with AI Presentation Generator"
    
    return slide

def create_title_content_slide(prs, section, points):
    """Create a slide with title and content"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = section
    
    text_frame = content_shape.text_frame
    text_frame.clear()
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.word_wrap = True
    
    first = True
    for point in points:
        level = point.get("level", 0)
        
        if first:
            p = text_frame.paragraphs[0]
            first = False
        else:
            p = text_frame.add_paragraph()
        
        p.text = point["text"]
        p.level = level
    
    return slide

def create_image_left_text_right_slide(prs, section, points, image_path):
    """Create a slide with image on left and text on right"""
    # Use a blank layout
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = section
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    
    # Add image on left
    img_left = Inches(0.5)
    img_top = Inches(1.5)
    img_width = Inches(4.5)
    img_height = Inches(5.5)
    slide.shapes.add_picture(image_path, img_left, img_top, width=img_width, height=img_height)
    
    # Add text on right
    text_left = Inches(5.5)
    text_top = Inches(1.5)
    text_width = Inches(4.5)
    text_height = Inches(5.5)
    text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    first = True
    for point in points:
        level = point.get("level", 0)
        
        if first:
            p = text_frame.paragraphs[0]
            first = False
        else:
            p = text_frame.add_paragraph()
        
        p.text = point["text"]
        p.level = level
        p.font.size = Pt(14)
    
    return slide

def create_image_right_text_left_slide(prs, section, points, image_path):
    """Create a slide with image on right and text on left"""
    # Use a blank layout
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = section
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    
    # Add text on left
    text_left = Inches(0.5)
    text_top = Inches(1.5)
    text_width = Inches(4.5)
    text_height = Inches(5.5)
    text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    first = True
    for point in points:
        level = point.get("level", 0)
        
        if first:
            p = text_frame.paragraphs[0]
            first = False
        else:
            p = text_frame.add_paragraph()
        
        p.text = point["text"]
        p.level = level
        p.font.size = Pt(14)
    
    # Add image on right
    img_left = Inches(5.5)
    img_top = Inches(1.5)
    img_width = Inches(4.5)
    img_height = Inches(5.5)
    slide.shapes.add_picture(image_path, img_left, img_top, width=img_width, height=img_height)
    
    return slide

def create_two_column_slide(prs, section, points):
    """Create a slide with two columns of text"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = section
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    
    # Split points into two columns
    mid_point = len(points) // 2
    left_points = points[:mid_point]
    right_points = points[mid_point:]
    
    # Add left column
    left_text_left = Inches(0.5)
    left_text_top = Inches(1.5)
    left_text_width = Inches(4.5)
    left_text_height = Inches(5.5)
    left_text_box = slide.shapes.add_textbox(left_text_left, left_text_top, left_text_width, left_text_height)
    left_text_frame = left_text_box.text_frame
    left_text_frame.word_wrap = True
    
    first = True
    for point in left_points:
        level = point.get("level", 0)
        
        if first:
            p = left_text_frame.paragraphs[0]
            first = False
        else:
            p = left_text_frame.add_paragraph()
        
        p.text = point["text"]
        p.level = level
        p.font.size = Pt(14)
    
    # Add right column
    right_text_left = Inches(5.5)
    right_text_top = Inches(1.5)
    right_text_width = Inches(4.5)
    right_text_height = Inches(5.5)
    right_text_box = slide.shapes.add_textbox(right_text_left, right_text_top, right_text_width, right_text_height)
    right_text_frame = right_text_box.text_frame
    right_text_frame.word_wrap = True
    
    first = True
    for point in right_points:
        level = point.get("level", 0)
        
        if first:
            p = right_text_frame.paragraphs[0]
            first = False
        else:
            p = right_text_frame.add_paragraph()
        
        p.text = point["text"]
        p.level = level
        p.font.size = Pt(14)
    
    return slide

def create_comparison_slide(prs, section, points):
    """Create a comparison slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = section
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    
    # Add left header
    left_header_left = Inches(0.5)
    left_header_top = Inches(1.5)
    left_header_width = Inches(4.5)
    left_header_height = Inches(0.5)
    left_header_box = slide.shapes.add_textbox(left_header_left, left_header_top, left_header_width, left_header_height)
    left_header_frame = left_header_box.text_frame
    left_header_frame.text = "Before / Option A"
    left_header_frame.paragraphs[0].font.size = Pt(18)
    left_header_frame.paragraphs[0].font.bold = True
    left_header_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add right header
    right_header_left = Inches(5.5)
    right_header_top = Inches(1.5)
    right_header_width = Inches(4.5)
    right_header_height = Inches(0.5)
    right_header_box = slide.shapes.add_textbox(right_header_left, right_header_top, right_header_width, right_header_height)
    right_header_frame = right_header_box.text_frame
    right_header_frame.text = "After / Option B"
    right_header_frame.paragraphs[0].font.size = Pt(18)
    right_header_frame.paragraphs[0].font.bold = True
    right_header_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add separator line
    line_left = Inches(5)
    line_top = Inches(1.5)
    line_width = Inches(0)
    line_height = Inches(5.5)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, line_left, line_top, line_width, line_height)
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line.line.fill.background()
    
    # Add left content
    left_text_left = Inches(0.5)
    left_text_top = Inches(2)
    left_text_width = Inches(4.5)
    left_text_height = Inches(5)
    left_text_box = slide.shapes.add_textbox(left_text_left, left_text_top, left_text_width, left_text_height)
    left_text_frame = left_text_box.text_frame
    left_text_frame.word_wrap = True
    
    # Add right content
    right_text_left = Inches(5.5)
    right_text_top = Inches(2)
    right_text_width = Inches(4.5)
    right_text_height = Inches(5)
    right_text_box = slide.shapes.add_textbox(right_text_left, right_text_top, right_text_width, right_text_height)
    right_text_frame = right_text_box.text_frame
    right_text_frame.word_wrap = True
    
    # Split points between left and right
    for i, point in enumerate(points):
        if i % 2 == 0 and len(left_text_frame.paragraphs) < 6:  # Limit to 6 points per side
            p = left_text_frame.add_paragraph() if i > 0 else left_text_frame.paragraphs[0]
            p.text = point["text"]
            p.font.size = Pt(14)
            p.level = 0
        elif len(right_text_frame.paragraphs) < 6:  # Limit to 6 points per side
            p = right_text_frame.add_paragraph() if i > 1 else right_text_frame.paragraphs[0]
            p.text = point["text"]
            p.font.size = Pt(14)
            p.level = 0
    
    return slide

def create_image_full_slide(prs, section, image_path):
    """Create a slide with a full-size image"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add image to fill the entire slide
    slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))
    
    # Add a semi-transparent title box at the bottom
    title_left = Inches(0)
    title_top = Inches(5.5)
    title_width = Inches(10)
    title_height = Inches(2)
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    
    # Add a semi-transparent background
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = RGBColor(0, 0, 0)
    title_box.fill.transparency = 0.3
    
    title_frame.text = section
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def create_conclusion_slide(prs, section, points):
    """Create a conclusion slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = section
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add content in the center
    content_left = Inches(1)
    content_top = Inches(2)
    content_width = Inches(8)
    content_height = Inches(4)
    content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Add key points
    for i, point in enumerate(points):
        if i < 3:  # Limit to 3 key points
            p = content_frame.add_paragraph() if i > 0 else content_frame.paragraphs[0]
            p.text = point["text"]
            p.font.size = Pt(20)
            p.level = 0
            p.alignment = PP_ALIGN.CENTER
    
    return slide

# Add a new function for the Thank You slide
def create_thank_you_slide(prs):
    """Create a separate Thank You slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add background color
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(25, 25, 112)  # Dark blue background
    background.line.fill.background()
    
    # Add Thank You text
    thank_you_left = Inches(1)
    thank_you_top = Inches(2.5)
    thank_you_width = Inches(8)
    thank_you_height = Inches(2)
    thank_you_box = slide.shapes.add_textbox(thank_you_left, thank_you_top, thank_you_width, thank_you_height)
    thank_you_frame = thank_you_box.text_frame
    thank_you_frame.text = "Thank You!"
    thank_you_frame.paragraphs[0].font.size = Pt(48)
    thank_you_frame.paragraphs[0].font.bold = True
    thank_you_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text
    thank_you_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add smaller subtitle
    subtitle_left = Inches(1)
    subtitle_top = Inches(5)
    subtitle_width = Inches(8)
    subtitle_height = Inches(1)
    subtitle_box = slide.shapes.add_textbox(subtitle_left, subtitle_top, subtitle_width, subtitle_height)
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Questions & Discussion"
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)  # Light gray text
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def list_available_themes(theme_folder="themes"):
    """List all available themes in the specified folder"""
    if not os.path.exists(theme_folder):
        os.makedirs(theme_folder)
        return []
    
    themes = [f for f in os.listdir(theme_folder) if f.endswith(".pptx")]
    return themes

def get_theme_path(theme_name, theme_folder="themes"):
    """Get the full path to a theme file"""
    return os.path.join(theme_folder, f"{theme_name}.pptx")


def create_presentation(outline, presentation_title, detail_level, filename="presentation.pptx", theme_path=None):
    """
    Creates a PowerPoint presentation from a structured outline.
    """
    try:
        # Load the theme from the specified path, if it exists.
        if theme_path and os.path.exists(theme_path):
            prs = Presentation(theme_path)
            
            # Delete all existing slides from the template.
            for i in range(len(prs.slides) - 1, -1, -1):
                rId = prs.slides._sldIdLst[i].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[i]
        else:
            # Create a blank presentation if no theme is found.
            prs = Presentation()
        
        # Add the main title slide.
        create_title_slide(prs, presentation_title)
        
        # Create a temporary directory for downloaded images to avoid clutter.
        with tempfile.TemporaryDirectory() as temp_dir:
            # Iterate through the outline to create each slide.
            for i, (section, points) in enumerate(outline.items()):
                slide_index = i + 1  # Offset by 1 for the title slide.
                
                # Determine the correct layout for the current slide.
                # The 'section' variable holds the slide title and is passed to the layout function.
                layout = determine_slide_layout(slide_index, detail_level, len(points), outline, section)
                
                print(f"Creating slide {slide_index}: {section} ({layout} layout)")
                
                # Create the slide based on the determined layout.
                if layout == "title_content":
                    create_title_content_slide(prs, section, points)
                
                elif layout in ["image_left_text_right", "image_right_text_left", "image_full"]:
                    # Search for relevant images for the slide.
                    queries = get_relevant_image_queries(presentation_title, section, points, detail_level == "detailed")
                    image_found = False
                    
                    for query in queries:
                        images = search_images(query, 3)
                        if images:
                            image_path = os.path.join(temp_dir, f"{section}_{images[0]['id']}.jpg")
                            if download_image(images[0]['download_url'], image_path):
                                image_found = True
                                if layout == "image_left_text_right":
                                    create_image_left_text_right_slide(prs, section, points, image_path)
                                elif layout == "image_right_text_left":
                                    create_image_right_text_left_slide(prs, section, points, image_path)
                                elif layout == "image_full":
                                    create_image_full_slide(prs, section, image_path)
                                break
                    
                    # Fallback to a content-only slide if no image is found.
                    if not image_found:
                        print(f"⚠️ No image found for '{section}', using title_content layout instead.")
                        create_title_content_slide(prs, section, points)
                
                elif layout == "two_column":
                    create_two_column_slide(prs, section, points)
                
                elif layout == "conclusion":
                    create_conclusion_slide(prs, section, points)
        
        # Add Thank You slide at the end
        create_thank_you_slide(prs)
        
        # Save the completed presentation to the specified filename.
        prs.save(filename)
        print(f"✅ Presentation saved as: {filename}")
        return filename
    
    except Exception as e:
        print(f"❌ Error creating presentation: {e}")
        # Re-raise the exception to be handled by the Flask app.
        raise

def ensure_conclusion_slide(outline, presentation_title):
    """Ensure the presentation always ends with a conclusion slide"""
    last_section = list(outline.keys())[-1] if outline else ""
    
    # If the last slide isn't a conclusion, add one
    if not any(keyword in last_section.lower() for keyword in ["conclusion", "summary", "wrap-up", "final"]):
        outline["Conclusion"] = [
            {"text": "Summary of key points and takeaways", "level": 0},
            {"text": "Future outlook and recommendations", "level": 0},
            {"text": "Q&A and discussion", "level": 0}
        ]
    return outline



def main():
    parser = argparse.ArgumentParser(description="Generate AI-powered presentations with varied layouts")
    parser.add_argument("--topic", type=str, default="Artificial Intelligence", help="Presentation topic")
    parser.add_argument("--detail", type=str, default="simple", choices=["simple", "detailed"], help="Detail level")
    parser.add_argument("--theme", type=str, help="Name of the theme file (without extension)")
    parser.add_argument("--output", type=str, default="presentation.pptx", help="Output filename")
    
    args = parser.parse_args()
    
    # List available themes
    theme_folder = "themes"
    available_themes = list_available_themes(theme_folder)
    
    print("🎯 AI Presentation Generator with Varied Layouts")
    print("=" * 50)
    print(f"Topic: {args.topic}")
    print(f"Detail Level: {args.detail}")
    
    # Handle theme selection
    if available_themes:
        print("Available Themes:")
        for theme in available_themes:
            print(f" - {os.path.splitext(theme)[0]}")
        
        if args.theme:
            theme_path = get_theme_path(args.theme, theme_folder)
            if not os.path.exists(theme_path):
                print(f"❌ Theme '{args.theme}' not found. Using default theme.")
                theme_path = None
        else:
            print("ℹ️  No theme specified. Using default theme.")
            theme_path = None
    else:
        print("ℹ️  No themes found in 'themes' folder. Using default theme.")
        theme_path = None
    
    print("=" * 50)
    
    # Generate outline
    outline, presentation_title = generate_outline(args.topic, args.detail)
    # Ensure conclusion slide
    outline = ensure_conclusion_slide(outline, presentation_title)
    # Create presentation with varied layouts
    create_presentation(outline, presentation_title, args.detail, args.output, theme_path)

    print("🎉 Presentation created successfully with varied layouts!")

if __name__ == "__main__":
    main()